from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()

# Define slide layout (Title and Content)
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]
blank_slide_layout = prs.slide_layouts[6]

# Helper function to add a slide with title and content
def add_slide(title, content_lines):
    slide = prs.slides.add_slide(content_slide_layout)
    title_placeholder = slide.shapes.title
    content_placeholder = slide.placeholders[1]
    title_placeholder.text = title
    content_placeholder.text = "\n".join(content_lines)

# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
slide.shapes.title.text = "Understanding Second Normal Form (2NF)"
slide.placeholders[1].text = "Database Normalization Principles\nPresenter: [Your Name]\nDate: [Insert Date]"

# Slide 2: What is 2NF?
add_slide("What is 2NF?", [
    "A table is in 2NF if:",
    "1. It is already in First Normal Form (1NF)",
    "2. It has no partial dependencies",
    "All non-key attributes must depend on the entire primary key"
])

# Slide 3: Key Concepts
add_slide("Key Terms You Should Know", [
    "- Primary Key: Unique identifier for a record",
    "- Composite Key: A key made from more than one column",
    "- Partial Dependency: Attribute depends on part of a composite key"
])

# Slide 4: Why Use 2NF?
add_slide("Importance of 2NF", [
    "✅ Eliminates redundant data",
    "✅ Prevents update, insert, and delete anomalies",
    "✅ Improves consistency and maintainability of data"
])

# Slide 5: Example of Table Not in 2NF
add_slide("Example - Table Before 2NF", [
    "| StudentID | CourseID | StudentName | CourseName | Grade |",
    "|-----------|----------|-------------|------------|-------|",
    "| 101       | C001     | Asha        | Math       | A     |",
    "| 101       | C002     | Asha        | Physics    | B+    |",
    "| 102       | C001     | John        | Math       | A-    |",
    "",
    "Problems:",
    "- StudentName depends only on StudentID",
    "- CourseName depends only on CourseID"
])

# Slide 6: Breaking the Table Into 2NF
add_slide("How to Normalize to 2NF", [
    "Split the original table into three tables:",
    "",
    "1. Student Table: | StudentID | StudentName |",
    "2. Course Table:  | CourseID  | CourseName  |",
    "3. Enrollment Table: | StudentID | CourseID | Grade |"
])

# Slide 7: Summary
add_slide("Recap: Second Normal Form (2NF)", [
    "- Removes partial dependencies",
    "- Every attribute depends on the full primary key",
    "- Results in cleaner, more reliable databases"
])

# Slide 8: Pro Tip
add_slide("Quick Tip for 2NF", [
    "📌 If your table has a single-column primary key and is already in 1NF,",
    "then it’s automatically in 2NF!"
])

# Slide 9: Questions & Practice
add_slide("Ready to Try It?", [
    "💬 Practice: Normalize this sample table",
    "📝 Ask questions if you're not sure",
    "👨🏽‍🏫 We’ll go through it together!"
])

# Save the presentation
output_path = "/mnt/data/Second_Normal_Form_2NF_Presentation.pptx"
prs.save(output_path)

output_path
